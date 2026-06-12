"""
management-overview.py
Generates management-overview.pptx — a high-level PowerPoint presentation
for senior management and business stakeholders.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x3A)
MID_BLUE    = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEB, 0xF7)
GREEN       = RGBColor(0x37, 0x8F, 0x5C)
AMBER       = RGBColor(0xF0, 0xA5, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY    = RGBColor(0x70, 0x70, 0x70)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
RED         = RGBColor(0xC0, 0x39, 0x2B)

OUTPUT_PATH = "/Volumes/projects/claudeCode/claudeProjects/release-manifest/docs/management-overview.pptx"

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width_pt=0):
    """Add a filled rectangle shape to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, font_size, bold, color,
             align=PP_ALIGN.LEFT, italic=False):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Dark blue bar at top, white title, optional light blue subtitle."""
    bar_height = 0.75 if subtitle else 0.65
    add_rect(slide, 0, 0, 13.33, bar_height, DARK_BLUE)
    add_text(slide, title, 0.35, 0.04, 12.5, 0.45,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.5, 12.5, 0.3,
                 font_size=12, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    return bar_height


def footer(slide):
    """Dark blue bottom strip with confidential label."""
    add_rect(slide, 0, 7.2, 13.33, 0.3, DARK_BLUE)
    add_text(slide, "Release Management Strategy  |  Confidential",
             0, 7.2, 13.33, 0.3,
             font_size=9, bold=False, color=MID_GREY, align=PP_ALIGN.CENTER)


def add_para_list(slide, items, left, top, width, height, font_size, color,
                  space_after_pt, bullet="• "):
    """Add a bulleted paragraph list using proper PPTX paragraphs."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(space_after_pt)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide1_cover(prs):
    """Slide 1 — Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark blue background
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)

    # Accent blue horizontal stripe at y=3.6"
    add_rect(slide, 0, 3.6, 13.33, 0.08, ACCENT_BLUE)

    # Large white title centered
    add_text(slide, "Release Management Strategy",
             0.5, 1.8, 12.33, 1.0,
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, "Controlled, Auditable, Bi-Weekly Delivery — Across All Teams",
             0.5, 2.8, 12.33, 0.65,
             font_size=20, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Below stripe tagline
    add_text(slide, "Bi-Weekly Sprint Cadence  ·  5 Teams  ·  Mixed Delivery Platforms",
             0.5, 3.85, 12.33, 0.5,
             font_size=15, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    footer(slide)


def slide2_challenge(prs):
    """Slide 2 — The Challenge We Were Facing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light grey background
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)

    header_bar(slide, "The Challenge")

    # LEFT column — Before
    lx = 0.35
    # Section label — dark blue bar
    add_rect(slide, lx, 0.85, 6.0, 0.4, DARK_BLUE)
    add_text(slide, "Before", lx + 0.1, 0.87, 5.8, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    problem_items = [
        "No single source of truth — who built what?",
        "Impossible to remove one team without rebuilding everything",
        "No record of exactly what was deployed to SIT vs UAT",
        "Teams on different tools (Git, SharePoint) — no unified process",
    ]
    card_top = 1.35
    for item in problem_items:
        # White card background
        add_rect(slide, lx, card_top, 5.8, 0.9, WHITE)
        # Red left border strip
        add_rect(slide, lx, card_top, 0.08, 0.9, RED)
        add_text(slide, item, lx + 0.18, card_top + 0.12, 5.5, 0.7,
                 font_size=13, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)
        card_top += 1.0

    # RIGHT column — After
    rx = 6.7
    add_rect(slide, rx, 0.85, 6.3, 0.4, GREEN)
    add_text(slide, "After", rx + 0.1, 0.87, 6.1, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    solution_items = [
        "One Excel manifest per sprint — CM owns every build decision",
        "Any team can be deferred or removed with a single checkbox change",
        "Every environment tagged to an exact version — always auditable",
        "Works for Git teams AND SharePoint teams in the same build",
    ]
    card_top = 1.35
    for item in solution_items:
        add_rect(slide, rx, card_top, 6.1, 0.9, WHITE)
        add_rect(slide, rx, card_top, 0.08, 0.9, GREEN)
        add_text(slide, item, rx + 0.18, card_top + 0.12, 5.8, 0.7,
                 font_size=13, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)
        card_top += 1.0

    footer(slide)


def slide3_five_teams(prs):
    """Slide 3 — How It Works: The Five Teams"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    hh = header_bar(slide, "Our Five Teams — Two Delivery Models")

    # Subtitle below header
    add_text(slide,
             "All teams feed into one controlled build — regardless of how they store their work",
             0.35, 0.8, 12.6, 0.4,
             font_size=12, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

    teams = [
        {
            "name": "Conversion",
            "title_color": AMBER,
            "badge": "SharePoint",
            "body": "Customer data\nETL pipelines\n\nDelivers via:\nSharePoint folder\n\nVersion tracked\nas artifact v1.x",
        },
        {
            "name": "Interfaces",
            "title_color": MID_BLUE,
            "badge": "Git",
            "body": "REST & file\nintegrations\n\nDelivers via:\nGit branch\n\nVersion tracked\nas commit SHA",
        },
        {
            "name": "Workflow Config",
            "title_color": MID_BLUE,
            "badge": "Git",
            "body": "Business workflow\nrules & approvals\n\nDelivers via:\nGit branch\n\nMay be deferred\nif not ready",
        },
        {
            "name": "Func Config",
            "title_color": MID_BLUE,
            "badge": "Git",
            "body": "Product config\n& region codes\n\nDelivers via:\nGit branch\n\nReference data\nfor all teams",
        },
        {
            "name": "COTS Platform",
            "title_color": DARK_BLUE,
            "badge": "Vendor",
            "body": "Third-party ETL\nplatform upgrade\n\nDelivers via:\nVendor package\n\nHotfixes tracked\nin build record",
        },
    ]

    card_w = 2.4
    card_h = 3.5
    gap = 0.2
    start_x = 0.35
    start_y = 1.35

    for i, team in enumerate(teams):
        cx = start_x + i * (card_w + gap)

        # Card white background
        add_rect(slide, cx, start_y, card_w, card_h, WHITE,
                 line_color=MID_GREY, line_width_pt=0.5)

        # Title bar
        add_rect(slide, cx, start_y, card_w, 0.5, team["title_color"])
        add_text(slide, team["name"], cx + 0.08, start_y + 0.05, card_w - 0.16, 0.4,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # Badge
        badge_color = AMBER if team["badge"] == "SharePoint" else (
            GREEN if team["badge"] == "Git" else ACCENT_BLUE)
        add_rect(slide, cx + card_w - 1.0, start_y + 0.55, 0.9, 0.28, badge_color)
        add_text(slide, team["badge"], cx + card_w - 1.0, start_y + 0.56, 0.9, 0.26,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Body text
        add_text(slide, team["body"], cx + 0.1, start_y + 0.92, card_w - 0.2, 2.5,
                 font_size=11, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

    # Bottom accent bar
    bar_y = start_y + card_h + 0.2
    add_rect(slide, 0.35, bar_y, 12.63, 0.45, ACCENT_BLUE)
    add_text(slide,
             "CM assembles all five into one controlled build — the manifest is the single record of what went in",
             0.35, bar_y + 0.05, 12.63, 0.35,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(slide)


def slide4_manifest(prs):
    """Slide 4 — The Build Manifest: Our Single Source of Truth"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header_bar(slide, "The Build Manifest — Excel File, CM-Owned")

    # ---- LEFT: Simulated Excel table ----
    lx = 0.35
    table_w = 5.5
    ty = 0.85

    # Table title bar
    add_rect(slide, lx, ty, table_w, 0.45, DARK_BLUE)
    add_text(slide, "manifest.xlsx — Sprint 12", lx + 0.1, ty + 0.06, table_w - 0.2, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Header row
    row_h = 0.45
    ty += 0.45
    add_rect(slide, lx, ty, table_w, row_h, DARK_GREY)
    headers = ["Team", "Source", "Included?", "Version", "ELM Items"]
    col_widths = [1.2, 0.9, 0.9, 1.1, 1.4]
    cx = lx
    for h, cw in zip(headers, col_widths):
        add_text(slide, h, cx + 0.04, ty + 0.08, cw - 0.08, row_h - 0.1,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        cx += cw

    # Data rows
    data_rows = [
        ("Conversion", "SharePoint", "✅ Yes", "v1.3", "ELM-1201, 1202, 1208", True),
        ("Interfaces", "Git", "✅ Yes", "b2c3d4e5", "ELM-1210, 1215", False),
        ("Workflow Config", "Git", "⛔ Deferred", "—", "ELM-1220, 1221", True),
        ("Func Config", "Git", "✅ Yes", "c3d4e5f6", "ELM-1230, 1231", False),
    ]

    for row_idx, (team, source, included, version, elms, is_alt) in enumerate(data_rows):
        ty += row_h
        row_bg = LIGHT_BLUE if is_alt else WHITE
        add_rect(slide, lx, ty, table_w, row_h, row_bg)

        # Included? cell colour
        inc_color = GREEN if "Yes" in included else AMBER

        cells = [team, source, included, version, elms]
        cx = lx
        for ci, (cell, cw) in enumerate(zip(cells, col_widths)):
            cell_color = DARK_GREY
            if ci == 2:
                cell_color = inc_color
            add_text(slide, cell, cx + 0.04, ty + 0.09, cw - 0.08, row_h - 0.12,
                     font_size=9, bold=(ci == 2), color=cell_color, align=PP_ALIGN.LEFT)
            cx += cw

    # ---- RIGHT: What the manifest controls ----
    rx = 6.2
    add_rect(slide, rx, 0.85, 6.8, 0.35, MID_BLUE)
    add_text(slide, "What the manifest controls", rx + 0.1, 0.87, 6.6, 0.3,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    bullet_items = [
        "Which teams are in this build — and which are held back",
        "The exact version (commit or file) deployed to each environment",
        "COTS platform version and hotfix list",
        "Who built it, when, and which ELM work items are included",
        "SIT and UAT deployment status — updated as environments progress",
    ]
    add_para_list(slide, bullet_items,
                  left=rx, top=1.3, width=6.8, height=2.8,
                  font_size=13, color=DARK_GREY, space_after_pt=10)

    # Callout box
    cy = 4.25
    add_rect(slide, rx, cy, 6.8, 1.5, LIGHT_BLUE, line_color=MID_BLUE, line_width_pt=1.5)
    add_text(slide, "CM approves every change to this file",
             rx + 0.15, cy + 0.1, 6.5, 0.4,
             font_size=13, bold=True, color=MID_BLUE, align=PP_ALIGN.LEFT)
    add_text(slide,
             "The manifest lives in a controlled repository. All edits are tracked with a full audit trail — who changed what and when.",
             rx + 0.15, cy + 0.5, 6.5, 0.9,
             font_size=12, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

    footer(slide)


def slide5_timeline(prs):
    """Slide 5 — The Sprint Lifecycle (Timeline)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header_bar(slide, "Sprint-12 in Numbers — From Start to Sign-Off")

    # Timeline baseline
    timeline_y = 3.4
    add_rect(slide, 0.5, timeline_y, 12.33, 0.08, ACCENT_BLUE)

    milestones = [
        {
            "label": "Day 1\nSprint Cut",
            "action": "Branches created\n(Git teams)",
        },
        {
            "label": "Days 2-9\nTeams Develop",
            "action": "Feature work\nin own systems",
        },
        {
            "label": "Day 10\nBuild Day",
            "action": "Manifest finalised\ninclude/defer decision",
        },
        {
            "label": "Day 10\nReport Generated",
            "action": "Release notes\nauto-generated",
        },
        {
            "label": "Day 11\nSIT Deploy",
            "action": "Tagged to exact\nversion deployed",
        },
        {
            "label": "Days 12-14\nSIT Testing",
            "action": "Defects raised\nin ELM",
        },
        {
            "label": "Day 15\nUAT Deploy",
            "action": "Business\nsign-off testing",
        },
        {
            "label": "Day 20\nSprint Close",
            "action": "Git branches\nmerged to main",
        },
    ]

    n = len(milestones)
    circle_d = 0.35
    total_w = 12.33
    spacing = total_w / (n - 1)
    start_x = 0.5

    for i, ms in enumerate(milestones):
        cx = start_x + i * spacing
        node_left = cx - circle_d / 2
        node_top = timeline_y - circle_d / 2 + 0.04

        # Circle node
        shape = slide.shapes.add_shape(
            9,  # oval
            Inches(node_left), Inches(node_top),
            Inches(circle_d), Inches(circle_d)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

        # Label below
        label_w = 1.4
        label_left = cx - label_w / 2
        add_text(slide, ms["label"], label_left, timeline_y + 0.28, label_w, 0.65,
                 font_size=10, bold=False, color=DARK_GREY, align=PP_ALIGN.CENTER)

        # Action above
        action_w = 1.5
        action_left = cx - action_w / 2
        add_text(slide, ms["action"], action_left, timeline_y - 0.95, action_w, 0.65,
                 font_size=9, bold=False, color=MID_BLUE, align=PP_ALIGN.CENTER)

    # Stat boxes below timeline
    stat_y = 5.0
    stat_data = [
        ("3 of 4 teams included", "1 team deferred — no impact on build"),
        ("1 COTS upgrade", "2 hotfixes tracked and documented"),
        ("100% auditable", "Every decision recorded in the manifest"),
    ]
    box_w = 4.0
    box_gap = 0.27
    for i, (title_t, body_t) in enumerate(stat_data):
        bx = 0.35 + i * (box_w + box_gap)
        add_rect(slide, bx, stat_y, box_w, 1.6, LIGHT_BLUE)
        add_text(slide, title_t, bx + 0.15, stat_y + 0.15, box_w - 0.3, 0.45,
                 font_size=14, bold=True, color=MID_BLUE, align=PP_ALIGN.LEFT)
        add_text(slide, body_t, bx + 0.15, stat_y + 0.65, box_w - 0.3, 0.8,
                 font_size=12, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

    footer(slide)


def slide6_risk(prs):
    """Slide 6 — Risk Control"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header_bar(slide, "Risk Control — We Can Remove Any Team at Any Time")

    card_top = 0.85
    card_h = 4.2

    # LEFT — Scenario 1
    lx = 0.35
    lw = 6.0

    add_rect(slide, lx, card_top, lw, 0.4, AMBER)
    add_text(slide, "Scenario 1: Team not ready on Build Day",
             lx + 0.1, card_top + 0.05, lw - 0.2, 0.32,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    sc1_top = card_top + 0.4
    sc1_h = card_h - 0.4
    add_rect(slide, lx, sc1_top, lw, sc1_h, WHITE)
    add_rect(slide, lx, sc1_top, 0.08, sc1_h, AMBER)

    add_text(slide, "Team deferred before build",
             lx + 0.18, sc1_top + 0.12, lw - 0.3, 0.4,
             font_size=14, bold=True, color=DARK_GREY, align=PP_ALIGN.LEFT)

    steps1 = [
        "CM marks team as 'Deferred' in the manifest Excel",
        "Their work stays safely in their branch or SharePoint folder",
        "Rest of the build proceeds unchanged",
        "No code surgery — no rebasing — no risk to other teams",
        "Deferred ELM items are noted and carried to the next sprint",
    ]
    add_para_list(slide, steps1,
                  left=lx + 0.18, top=sc1_top + 0.58,
                  width=lw - 0.3, height=2.6,
                  font_size=12, color=DARK_GREY, space_after_pt=8)

    # Outcome box — green
    out1_y = sc1_top + sc1_h - 0.65
    add_rect(slide, lx, out1_y, lw, 0.55, GREEN)
    add_text(slide, "Zero impact on the rest of the build",
             lx + 0.1, out1_y + 0.1, lw - 0.2, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # RIGHT — Scenario 2
    rx = 6.7
    rw = 6.3

    add_rect(slide, rx, card_top, rw, 0.4, RED)
    add_text(slide, "Scenario 2: Problem found in SIT testing",
             rx + 0.1, card_top + 0.05, rw - 0.2, 0.32,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    sc2_top = card_top + 0.4
    sc2_h = card_h - 0.4
    add_rect(slide, rx, sc2_top, rw, sc2_h, WHITE)
    add_rect(slide, rx, sc2_top, 0.08, sc2_h, RED)

    add_text(slide, "Team removed mid-SIT",
             rx + 0.18, sc2_top + 0.12, rw - 0.3, 0.4,
             font_size=14, bold=True, color=DARK_GREY, align=PP_ALIGN.LEFT)

    steps2 = [
        "CM creates a 'patch' manifest — a copy with that team excluded",
        "Patch build is deployed to SIT (other teams unaffected)",
        "Original manifest is preserved as audit record",
        "Issue is logged in ELM and tracked for next sprint",
        "No production risk — defect caught in SIT",
    ]
    add_para_list(slide, steps2,
                  left=rx + 0.18, top=sc2_top + 0.58,
                  width=rw - 0.3, height=2.6,
                  font_size=12, color=DARK_GREY, space_after_pt=8)

    # Outcome box — accent blue
    out2_y = sc2_top + sc2_h - 0.65
    add_rect(slide, rx, out2_y, rw, 0.55, ACCENT_BLUE)
    add_text(slide, "Original audit trail is never overwritten",
             rx + 0.1, out2_y + 0.1, rw - 0.2, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(slide)


def slide7_benefits(prs):
    """Slide 7 — Benefits Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark blue full background
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)

    # Accent blue stripe
    add_rect(slide, 0, 1.5, 13.33, 0.06, ACCENT_BLUE)

    # Title
    add_text(slide, "Why This Approach Works",
             0.5, 0.5, 12.33, 0.85,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, "Designed for compliance, flexibility, and scale",
             0.5, 1.15, 12.33, 0.4,
             font_size=16, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    tiles = [
        ("Full Audit Trail", "Every build decision is recorded with who, what, and when"),
        ("Works for Any Team", "Git, SharePoint, or vendor packages — all in one build"),
        ("No Build Held Hostage", "One team's issues never delay everyone else's release"),
        ("ELM Traceability", "Every change links back to an ELM work item end-to-end"),
        ("Repeatable Process", "Same scripts, same Excel template, every sprint — no surprises"),
        ("Compliance Ready", "Immutable audit log meets regulatory and CM board requirements"),
    ]

    tile_w = 4.0
    tile_h = 1.35
    gap_x = 0.27
    gap_y = 0.2
    start_x = 0.35
    start_y = 1.75
    accent_bar_w = 0.08

    for i, (title_t, body_t) in enumerate(tiles):
        row = i // 3
        col = i % 3
        tx = start_x + col * (tile_w + gap_x)
        ty = start_y + row * (tile_h + gap_y)

        add_rect(slide, tx, ty, tile_w, tile_h, WHITE)
        add_rect(slide, tx, ty, accent_bar_w, tile_h, MID_BLUE)

        add_text(slide, title_t,
                 tx + accent_bar_w + 0.12, ty + 0.1, tile_w - accent_bar_w - 0.18, 0.42,
                 font_size=14, bold=True, color=DARK_GREY, align=PP_ALIGN.LEFT)
        add_text(slide, body_t,
                 tx + accent_bar_w + 0.12, ty + 0.52, tile_w - accent_bar_w - 0.18, 0.72,
                 font_size=11, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

    footer(slide)


def slide8_cta(prs):
    """Slide 8 — Summary / Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark blue full background
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)

    # Accent blue stripe
    add_rect(slide, 0, 2.0, 13.33, 0.06, ACCENT_BLUE)

    # Title
    add_text(slide, "This Strategy Is Live and Proven",
             0.5, 0.85, 12.33, 0.9,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, "Sprint-12 completed end-to-end using this process",
             0.5, 1.6, 12.33, 0.45,
             font_size=16, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Three summary boxes
    boxes = [
        ("5 Teams Onboarded", "Mix of Git and SharePoint — all managed from one manifest"),
        ("Bi-Weekly Cadence", "Sprint cut → SIT → UAT → Release in under 3 weeks"),
        ("Zero Process Gaps", "Deferred teams, mid-SIT removals, COTS upgrades — all covered"),
    ]
    box_w = 4.1
    box_h = 2.2
    box_gap = 0.27
    box_y = 2.25

    for i, (title_t, body_t) in enumerate(boxes):
        bx = 0.35 + i * (box_w + box_gap)
        add_rect(slide, bx, box_y, box_w, box_h, MID_BLUE)
        add_text(slide, title_t,
                 bx + 0.2, box_y + 0.2, box_w - 0.4, 0.55,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, body_t,
                 bx + 0.2, box_y + 0.82, box_w - 0.4, 1.2,
                 font_size=13, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

    # Wide callout box — slightly lighter dark blue, accent blue left border
    callout_y = 4.7
    callout_fill = RGBColor(0x26, 0x35, 0x55)
    add_rect(slide, 0.35, callout_y, 12.63, 1.55, callout_fill)
    add_rect(slide, 0.35, callout_y, 0.1, 1.55, ACCENT_BLUE)
    add_text(slide,
             "The manifest Excel file is the contract between the CM team and every delivery team. "
             "It is the only document that needs to be reviewed before any build goes to SIT or UAT.",
             0.65, callout_y + 0.25, 12.1, 1.1,
             font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    footer(slide)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide1_cover(prs)
    slide2_challenge(prs)
    slide3_five_teams(prs)
    slide4_manifest(prs)
    slide5_timeline(prs)
    slide6_risk(prs)
    slide7_benefits(prs)
    slide8_cta(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
