#!/usr/bin/env python3
"""
generate_pptx.py
Generates sprint-12-demo.pptx — a slide deck for the release management demo.
Run from the repo root:  python3 docs/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "sprint-12-demo.pptx")

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x3A)   # slide background / header bar
MID_BLUE    = RGBColor(0x1F, 0x4E, 0x79)   # section accent
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)   # highlight / badge
LIGHT_BLUE  = RGBColor(0xDE, 0xEB, 0xF7)   # card background
GREEN       = RGBColor(0x37, 0x8F, 0x5C)   # included / success
AMBER       = RGBColor(0xF0, 0xA5, 0x00)   # deferred / warning
RED         = RGBColor(0xC0, 0x39, 0x2B)   # blocked / error
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY    = RGBColor(0x70, 0x70, 0x70)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Dark blue full-width header bar at the top."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=DARK_BLUE)
    add_text_box(slide, title,
                 Inches(0.35), Inches(0.12), Inches(12), Inches(0.6),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.72), Inches(12), Inches(0.35),
                     font_size=14, color=RGBColor(0xBD, 0xD7, 0xEE))


def section_label(slide, text, left, top, width=Inches(3), color=ACCENT_BLUE):
    add_rect(slide, left, top, width, Inches(0.32), fill_color=color)
    add_text_box(slide, text, left + Inches(0.1), top + Inches(0.03),
                 width - Inches(0.2), Inches(0.28),
                 font_size=11, bold=True, color=WHITE)


def card(slide, left, top, width, height,
         title, body_lines, title_color=MID_BLUE, badge=None, badge_color=GREEN):
    add_rect(slide, left, top, width, height, fill_color=LIGHT_BLUE,
             line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
    # title bar inside card
    add_rect(slide, left, top, width, Inches(0.35), fill_color=title_color)
    title_text = f"{title}  {badge}" if badge else title
    add_text_box(slide, title_text,
                 left + Inches(0.1), top + Inches(0.04),
                 width - Inches(0.2), Inches(0.28),
                 font_size=11, bold=True, color=WHITE)
    # body
    body = "\n".join(body_lines)
    add_text_box(slide, body,
                 left + Inches(0.12), top + Inches(0.4),
                 width - Inches(0.24), height - Inches(0.5),
                 font_size=10.5, color=DARK_GREY, wrap=True)


def bullet_list(slide, items, left, top, width, height,
                font_size=13, color=DARK_GREY, bullet="▸  "):
    text = "\n".join(f"{bullet}{item}" for item in items)
    add_text_box(slide, text, left, top, width, height,
                 font_size=font_size, color=color, wrap=True)


def phase_badge(slide, number, title, left, top):
    """Numbered phase circle + title."""
    # circle
    circ = slide.shapes.add_shape(9, left, top, Inches(0.48), Inches(0.48))  # oval
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_BLUE
    circ.line.fill.background()
    add_text_box(slide, str(number),
                 left + Inches(0.07), top + Inches(0.05),
                 Inches(0.34), Inches(0.38),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 left + Inches(0.56), top + Inches(0.07),
                 Inches(3.2), Inches(0.36),
                 font_size=13, bold=True, color=DARK_BLUE)


def footer(slide, text="Sprint-12 Release Management Demo  |  Git Branching Strategy Prototype"):
    add_rect(slide, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.28),
             fill_color=DARK_BLUE)
    add_text_box(slide, text,
                 Inches(0.25), SLIDE_H - Inches(0.26),
                 SLIDE_W - Inches(0.5), Inches(0.24),
                 font_size=9, color=RGBColor(0x9D, 0xB8, 0xD2),
                 align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# full background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BLUE)
# accent stripe
add_rect(slide, 0, Inches(3.6), SLIDE_W, Inches(0.06), fill_color=ACCENT_BLUE)

add_text_box(slide, "Release Management",
             Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Git Branching Strategy — Sprint-12 End-to-End Demo",
             Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.65),
             font_size=22, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)

add_text_box(slide, "Sprint-12  ·  9 Jun – 20 Jun 2026",
             Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.4),
             font_size=14, color=RGBColor(0x9D, 0xB8, 0xD2), align=PP_ALIGN.CENTER)

add_text_box(slide, "5 Teams  ·  Polyrepo  ·  IBM ELM Integrated",
             Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.4),
             font_size=14, color=RGBColor(0x9D, 0xB8, 0xD2), align=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview: Strategy at a Glance
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Strategy at a Glance",
           "Polyrepo + lightweight Release Manifest Repo as CM coordination hub")

# Three pillar cards
pillars = [
    ("Polyrepo Structure",
     ["5 independent team repos", "No cross-team merge conflicts", "Each team owns their pipeline"]),
    ("Manifest-Driven Builds",
     ["manifest.yml = single source of truth", "CM decides include/defer per team", "Auditable via git history"]),
    ("Environment Tags",
     ["sit/sprint-NN  uat/sprint-NN", "release/sprint-NN on main", "Exact commit pinned forever"]),
]
for i, (title, lines) in enumerate(pillars):
    left = Inches(0.35 + i * 4.32)
    card(slide, left, Inches(1.4), Inches(4.0), Inches(2.2),
         title, lines, title_color=MID_BLUE)

# Bottom: repo list
add_rect(slide, Inches(0.35), Inches(3.85), Inches(12.6), Inches(0.35), fill_color=MID_BLUE)
add_text_box(slide, "Repositories in Play",
             Inches(0.45), Inches(3.88), Inches(12.0), Inches(0.28),
             font_size=12, bold=True, color=WHITE)

repos = [
    ("repo-release-manifest", "CM hub — manifests, scripts, runbook", ACCENT_BLUE),
    ("repo-conversion",       "Pentaho ETL — customer data pipelines", GREEN),
    ("repo-interfaces",       "REST integrations, Pentaho jobs",       GREEN),
    ("repo-workflow-config",  "Workflow rules, approval chains",        AMBER),
    ("repo-func-config",      "Region codes, product hierarchy",        GREEN),
]
for i, (name, desc, color) in enumerate(repos):
    col = i % 3
    row = i // 3
    lft = Inches(0.35 + col * 4.32)
    tp  = Inches(4.35 + row * 1.05)
    add_rect(slide, lft, tp, Inches(4.0), Inches(0.85),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_rect(slide, lft, tp, Inches(0.07), Inches(0.85), fill_color=color)
    add_text_box(slide, name, lft + Inches(0.15), tp + Inches(0.06),
                 Inches(3.7), Inches(0.32), font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide, desc, lft + Inches(0.15), tp + Inches(0.44),
                 Inches(3.7), Inches(0.28), font_size=10, color=MID_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Sprint Lifecycle Overview (phases 1-8)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Sprint-12 Lifecycle — 8 Phases",
           "From sprint cut to UAT close in one repeatable process")

phases = [
    (1, "Cut Sprint",        "CM creates sprint/sprint-12\nbranches in all team repos"),
    (2, "Teams Develop",     "Feature branches off sprint/sprint-12\nELM-XXXX: commit convention"),
    (3, "Build Day Manifest","CM fills manifest.yml\nincluded / deferred decision"),
    (4, "Generate Report",   "generate-report.sh sprint-12\nproduces RELEASE-NOTES.md"),
    (5, "Deploy to SIT",     "Included teams deployed\ntag-environments.sh sprint-12 sit"),
    (6, "SIT Testing",       "Optional: patch manifest\nif team removed mid-SIT"),
    (7, "Deploy to UAT",     "tag-environments.sh sprint-12 uat\nuat/sprint-12 tags created"),
    (8, "Close Sprint",      "close-sprint.sh sprint-12\nMerge to main, release tags"),
]

cols = 4
box_w = Inches(3.0)
box_h = Inches(1.55)
gap_x = Inches(0.25)
gap_y = Inches(0.3)
start_x = Inches(0.35)
start_y = Inches(1.35)

for i, (num, title, desc) in enumerate(phases):
    col = i % cols
    row = i // cols
    lft = start_x + col * (box_w + gap_x)
    tp  = start_y + row * (box_h + gap_y)

    color = ACCENT_BLUE if num not in (6,) else AMBER
    add_rect(slide, lft, tp, box_w, box_h,
             fill_color=WHITE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
    # phase number circle
    circ = slide.shapes.add_shape(9, lft + Inches(0.1), tp + Inches(0.1),
                                  Inches(0.38), Inches(0.38))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_text_box(slide, str(num),
                 lft + Inches(0.1), tp + Inches(0.1),
                 Inches(0.38), Inches(0.38),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 lft + Inches(0.55), tp + Inches(0.12),
                 Inches(2.3), Inches(0.32),
                 font_size=12, bold=True, color=DARK_BLUE)
    add_text_box(slide, desc,
                 lft + Inches(0.12), tp + Inches(0.55),
                 box_w - Inches(0.24), Inches(0.85),
                 font_size=10, color=MID_GREY, wrap=True)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Phase 1: Cut the Sprint
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 1 — Cut the Sprint",
           "CM runs cut-sprint.sh on Day 1 to create sprint branches and scaffold the manifest")

# Left: what the script does
section_label(slide, "What the script does", Inches(0.35), Inches(1.35))
steps = [
    "Fetches latest main in each team repo",
    "Creates sprint/sprint-12 branch from main",
    "Pushes the branch to GitHub",
    "Scaffolds releases/sprint-12/manifest.yml from template",
    "Scaffolds releases/sprint-12/cots-sprint-12.md",
    "Creates placeholder RELEASE-NOTES.md",
]
bullet_list(slide, steps, Inches(0.35), Inches(1.75), Inches(5.8), Inches(3.5))

# Right: command + output
section_label(slide, "Command", Inches(6.5), Inches(1.35), width=Inches(1.5))
add_rect(slide, Inches(6.5), Inches(1.75), Inches(6.5), Inches(0.7),
         fill_color=DARK_BLUE)
add_text_box(slide, "$ ./scripts/cut-sprint.sh sprint-12",
             Inches(6.6), Inches(1.8), Inches(6.3), Inches(0.55),
             font_size=13, color=RGBColor(0x9C, 0xE8, 0x7A))

section_label(slide, "Output — branches created in GitHub", Inches(6.5), Inches(2.6),
              width=Inches(3.8))
branch_rows = [
    ("repo-conversion",       "sprint/sprint-12  ✅"),
    ("repo-interfaces",       "sprint/sprint-12  ✅"),
    ("repo-workflow-config",  "sprint/sprint-12  ✅"),
    ("repo-func-config",      "sprint/sprint-12  ✅"),
]
for i, (repo, branch) in enumerate(branch_rows):
    tp = Inches(3.05 + i * 0.65)
    add_rect(slide, Inches(6.5), tp, Inches(6.5), Inches(0.55),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_text_box(slide, repo,  Inches(6.65), tp + Inches(0.08),
                 Inches(3.5), Inches(0.38), font_size=11, bold=True, color=MID_BLUE)
    add_text_box(slide, branch, Inches(10.2), tp + Inches(0.08),
                 Inches(2.5), Inches(0.38), font_size=11, color=GREEN)

# Teams also notified
section_label(slide, "Teams are notified of branch name", Inches(0.35), Inches(5.4),
              width=Inches(4.2))
add_text_box(slide, "sprint/sprint-12 — create your feature branches off this",
             Inches(0.35), Inches(5.8), Inches(5.8), Inches(0.5),
             font_size=12, italic=True, color=MID_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Phase 2: Teams Develop
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 2 — Teams Develop in Their Own Repos",
           "Each team works independently on feature/ELM-XXXX-* branches off sprint/sprint-12")

teams_work = [
    ("Conversion Team ✅",         MID_BLUE,  [
        "ELM-1201  customer_address_mapping.ktr v1.2",
        "  → suburb normalisation + PO Box detection",
        "ELM-1202  account_load.sql null guards",
        "  → reject null customer_id / empty account_type",
        "ELM-1208  rename_audit_column.sql (idempotent)",
        "  → created_by → audit_user for COTS 2.4.0",
    ]),
    ("Interfaces Team ✅",         MID_BLUE,  [
        "ELM-1210  rest_endpoints.properties",
        "  → all URLs moved to /v2, timeout 30s → 60s",
        "ELM-1215  retry_config.properties",
        "  → retry on 408 / 503 / 504 status codes",
        "  → circuit breaker added (10 failures / 60s)",
        "",
    ]),
    ("Workflow Config  ⛔ DEFERRED", AMBER,   [
        "ELM-1221  workflow_rules.xml v1.2  ✅ done",
        "  → Account closure workflow rule added",
        "ELM-1220  dynamic approval chains  ❌ blocked",
        "  → Pending business sign-off",
        "  → Too risky to deploy partial change",
        "  → CM decision: defer whole team",
    ]),
    ("Func Config Team ✅",         MID_BLUE,  [
        "ELM-1230  region_codes.sql",
        "  → MED, MEA, NAC regions added",
        "ELM-1231  product_hierarchy.sql",
        "  → Mediterranean & MEA product variants",
        "  → Balance thresholds updated for new markets",
        "",
    ]),
]

col_w = Inches(3.0)
col_h = Inches(4.3)
gap   = Inches(0.26)
start = Inches(0.35)

for i, (title, tc, lines) in enumerate(teams_work):
    lft = start + i * (col_w + gap)
    card(slide, lft, Inches(1.35), col_w, col_h, title, lines, title_color=tc)

# Commit convention callout at bottom
add_rect(slide, Inches(0.35), Inches(5.9), Inches(12.6), Inches(0.45), fill_color=DARK_BLUE)
add_text_box(slide,
             "Commit convention:  ELM-XXXX: description  — links every change to IBM ELM work item",
             Inches(0.5), Inches(5.93), Inches(12.0), Inches(0.38),
             font_size=12, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
             align=PP_ALIGN.CENTER)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Phase 3: Build Day Manifest
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 3 — Build Day: Finalise the Manifest",
           "CM opens manifest.yml and makes the include/defer decision for each team")

# Manifest YAML preview (left panel)
section_label(slide, "releases/sprint-12/manifest.yml", Inches(0.35), Inches(1.35),
              width=Inches(4.5))
yaml_preview = (
    "release: sprint-12\n"
    "build_date: \"2026-06-20\"\n"
    "built_by: john.smith\n\n"
    "teams:\n"
    "  conversion:\n"
    "    included: true\n"
    "    commit: a1b2c3d4\n"
    "    elm_work_items: [ELM-1201, ELM-1202, ELM-1208]\n\n"
    "  workflow_config:\n"
    "    included: false\n"
    "    reason: \"ELM-1220 blocked — deferred to sprint-13\"\n\n"
    "cots:\n"
    "  version: \"2.4.0\"\n"
    "  previous_version: \"2.3.8\"\n"
    "  hotfixes_included: [HF-2024-011, HF-2024-012]"
)
add_rect(slide, Inches(0.35), Inches(1.75), Inches(5.8), Inches(4.5), fill_color=DARK_BLUE)
add_text_box(slide, yaml_preview,
             Inches(0.5), Inches(1.85), Inches(5.5), Inches(4.3),
             font_size=10, color=RGBColor(0xAB, 0xDE, 0x78))

# Right: decision table
section_label(slide, "Build Day Decisions", Inches(6.5), Inches(1.35), width=Inches(3.0))
decisions = [
    ("Conversion",      "✅ Included",  GREEN, "All 3 ELM items complete"),
    ("Interfaces",      "✅ Included",  GREEN, "v2 upgrade + retry ready"),
    ("Workflow Config", "⛔ Deferred",  AMBER, "ELM-1220 blocked"),
    ("Func Config",     "✅ Included",  GREEN, "Region data complete"),
]
for i, (team, status, color, note) in enumerate(decisions):
    tp = Inches(1.8 + i * 0.85)
    add_rect(slide, Inches(6.5), tp, Inches(6.5), Inches(0.72),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_rect(slide, Inches(6.5), tp, Inches(0.07), Inches(0.72), fill_color=color)
    add_text_box(slide, team,   Inches(6.65), tp + Inches(0.06),
                 Inches(2.8), Inches(0.28), font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide, status, Inches(9.5),  tp + Inches(0.06),
                 Inches(1.5), Inches(0.28), font_size=11, bold=True, color=color)
    add_text_box(slide, note,   Inches(6.65), tp + Inches(0.38),
                 Inches(4.2), Inches(0.25), font_size=10, color=MID_GREY)

# COTS callout
add_rect(slide, Inches(6.5), Inches(5.35), Inches(6.5), Inches(0.85),
         fill_color=LIGHT_BLUE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
add_text_box(slide, "COTS: Acme ETL Platform 2.3.8 → 2.4.0",
             Inches(6.65), Inches(5.42), Inches(6.0), Inches(0.32),
             font_size=12, bold=True, color=MID_BLUE)
add_text_box(slide, "HF-2024-011, HF-2024-012  ·  Requires ELM-1208 migration (audit_user rename)",
             Inches(6.65), Inches(5.75), Inches(6.0), Inches(0.32),
             font_size=10, color=MID_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Phase 4: Generate Release Notes
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 4 — Generate the Release Notes",
           "./scripts/generate-report.sh sprint-12  →  releases/sprint-12/RELEASE-NOTES.md")

section_label(slide, "Command", Inches(0.35), Inches(1.35), width=Inches(1.5))
add_rect(slide, Inches(0.35), Inches(1.75), Inches(5.8), Inches(0.65), fill_color=DARK_BLUE)
add_text_box(slide, "$ ./scripts/generate-report.sh sprint-12",
             Inches(0.5), Inches(1.82), Inches(5.5), Inches(0.5),
             font_size=13, color=RGBColor(0x9C, 0xE8, 0x7A))

section_label(slide, "RELEASE-NOTES.md structure", Inches(0.35), Inches(2.6), width=Inches(4.0))
sections = [
    "# Release Notes — sprint-12",
    "| Sprint dates | 2026-06-09 → 2026-06-20 |",
    "| Built by     | john.smith              |",
    "",
    "## Teams Included in This Build",
    "### Conversion ✅",
    "- Branch: sprint/sprint-12   Commit: a1b2c3d4",
    "- ELM items: ELM-1201, ELM-1202, ELM-1208",
    "",
    "## Teams Deferred (not in this build)",
    "### Workflow Config ⏸",
    "- Reason: ELM-1220 blocked pending sign-off",
    "",
    "## COTS Product",
    "| Version | 2.4.0  Previous: 2.3.8 |",
    "| Hotfixes | HF-2024-011, HF-2024-012 |",
    "",
    "## Environment Status",
    "| SIT | deployed | 2026-06-20 | — |",
    "| UAT | pending  | —          | — |",
]
add_rect(slide, Inches(0.35), Inches(3.0), Inches(5.8), Inches(3.65), fill_color=WHITE,
         line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
add_text_box(slide, "\n".join(sections),
             Inches(0.5), Inches(3.1), Inches(5.5), Inches(3.4),
             font_size=9, color=DARK_GREY)

# Right: what the report is used for
section_label(slide, "How this report is used", Inches(6.5), Inches(1.35), width=Inches(3.5))
uses = [
    "Committed to repo-release-manifest",
    "PR raised for CM team review (CODEOWNERS)",
    "Shared with stakeholders as official build record",
    "Audit trail: who built what, from which commit",
    "Deferred teams logged — tracked to next sprint",
    "COTS hotfixes documented for compliance",
]
bullet_list(slide, uses, Inches(6.5), Inches(1.75), Inches(6.5), Inches(3.0), font_size=13)

add_rect(slide, Inches(6.5), Inches(5.0), Inches(6.5), Inches(1.2),
         fill_color=LIGHT_BLUE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
add_text_box(slide, "Implementation: Python + PyYAML",
             Inches(6.65), Inches(5.08), Inches(6.0), Inches(0.32),
             font_size=12, bold=True, color=MID_BLUE)
add_text_box(slide, (
    "Shell YAML parsing (grep/awk) failed on inline arrays [ELM-1201, ELM-1202].\n"
    "Rewritten as generate-report.py — robust, handles all manifest structures."
    ),
             Inches(6.65), Inches(5.45), Inches(6.0), Inches(0.65),
             font_size=10, color=MID_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Phase 5: Deploy to SIT and Tag
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 5 — Deploy to SIT and Tag",
           "After deployment: ./scripts/tag-environments.sh sprint-12 sit")

section_label(slide, "What gets tagged", Inches(0.35), Inches(1.35), width=Inches(2.8))
tag_rows = [
    ("repo-conversion",      "sit/sprint-12 → commit a1b2c3d4", GREEN,  True),
    ("repo-interfaces",      "sit/sprint-12 → commit b2c3d4e5", GREEN,  True),
    ("repo-func-config",     "sit/sprint-12 → commit c3d4e5f6", GREEN,  True),
    ("repo-workflow-config", "(skipped — included: false)",      AMBER,  False),
]
for i, (repo, tag, color, tagged) in enumerate(tag_rows):
    tp = Inches(1.75 + i * 0.85)
    add_rect(slide, Inches(0.35), tp, Inches(6.2), Inches(0.72),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_rect(slide, Inches(0.35), tp, Inches(0.07), Inches(0.72), fill_color=color)
    add_text_box(slide, repo, Inches(0.5), tp + Inches(0.06),
                 Inches(3.0), Inches(0.28), font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide, tag, Inches(0.5), tp + Inches(0.38),
                 Inches(5.8), Inches(0.28), font_size=10,
                 color=GREEN if tagged else AMBER)

# Right: what the tag proves
section_label(slide, "Why tags matter", Inches(7.0), Inches(1.35), width=Inches(2.5))
reasons = [
    "Survives branch deletion after sprint close",
    "Tells ops exactly what code is in each env",
    "Auditable: tag message records who deployed when",
    "Enables rollback: redeploy from tag at any time",
    "Patch releases get their own tag (sprint-12-patch1)",
]
bullet_list(slide, reasons, Inches(7.0), Inches(1.75), Inches(6.0), Inches(2.5), font_size=13)

# manifest update
section_label(slide, "manifest.yml updated automatically", Inches(7.0), Inches(4.4),
              width=Inches(4.2))
add_rect(slide, Inches(7.0), Inches(4.8), Inches(6.0), Inches(1.4), fill_color=DARK_BLUE)
add_text_box(slide,
             "environments:\n"
             "  sit:\n"
             "    status: deployed\n"
             "    deployed_by: john.smith\n"
             "    deployed_date: 2026-06-20 14:30",
             Inches(7.15), Inches(4.88), Inches(5.7), Inches(1.25),
             font_size=11, color=RGBColor(0xAB, 0xDE, 0x78))

# Mid-SIT scenario box
add_rect(slide, Inches(0.35), Inches(5.55), Inches(6.2), Inches(1.0),
         fill_color=RGBColor(0xFF, 0xF3, 0xCD), line_color=AMBER, line_width_pt=1.0)
add_text_box(slide, "⚠  Mid-SIT Issue?",
             Inches(0.5), Inches(5.6), Inches(5.8), Inches(0.32),
             font_size=12, bold=True, color=AMBER)
add_text_box(slide, "Create sprint-12-patch1/manifest.yml  →  set team to included: false  →  re-tag",
             Inches(0.5), Inches(5.95), Inches(5.8), Inches(0.45),
             font_size=10, color=DARK_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Phase 8: Close the Sprint
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Phase 8 — Close the Sprint (UAT Passed)",
           "./scripts/close-sprint.sh sprint-12  —  merge to main, tag releases, carry over deferrals")

# Left: included teams flow
section_label(slide, "Included teams — merged to main", Inches(0.35), Inches(1.35),
              width=Inches(4.0))
merged = [
    ("repo-conversion",  "sprint/sprint-12 → main  +  release/sprint-12"),
    ("repo-interfaces",  "sprint/sprint-12 → main  +  release/sprint-12"),
    ("repo-func-config", "sprint/sprint-12 → main  +  release/sprint-12"),
]
for i, (repo, action) in enumerate(merged):
    tp = Inches(1.8 + i * 0.9)
    add_rect(slide, Inches(0.35), tp, Inches(6.0), Inches(0.72),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_rect(slide, Inches(0.35), tp, Inches(0.07), Inches(0.72), fill_color=GREEN)
    add_text_box(slide, repo,   Inches(0.5), tp + Inches(0.06),
                 Inches(3.0), Inches(0.28), font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide, action, Inches(0.5), tp + Inches(0.38),
                 Inches(5.6), Inches(0.28), font_size=10, color=GREEN)

# Deferred carryover
section_label(slide, "Deferred team — sprint branch stays open", Inches(0.35), Inches(4.7),
              width=Inches(4.8))
add_rect(slide, Inches(0.35), Inches(5.1), Inches(6.0), Inches(1.0),
         fill_color=RGBColor(0xFF, 0xF3, 0xCD), line_color=AMBER, line_width_pt=1.0)
add_text_box(slide, "repo-workflow-config",
             Inches(0.5), Inches(5.18), Inches(5.6), Inches(0.28),
             font_size=11, bold=True, color=AMBER)
add_text_box(slide, "sprint/sprint-12 stays open  →  ELM-1220 continues in sprint-13 manifest",
             Inches(0.5), Inches(5.5), Inches(5.6), Inches(0.45),
             font_size=10, color=DARK_GREY)

# Right: tag hierarchy
section_label(slide, "Complete tag history after close", Inches(6.7), Inches(1.35),
              width=Inches(4.0))
tag_levels = [
    ("sit/sprint-12",     "Pinpoints SIT-tested commit",    ACCENT_BLUE),
    ("uat/sprint-12",     "Pinpoints UAT-approved commit",  MID_BLUE),
    ("release/sprint-12", "Production-ready main snapshot", GREEN),
]
for i, (tag, desc, color) in enumerate(tag_levels):
    tp = Inches(1.8 + i * 1.15)
    add_rect(slide, Inches(6.7), tp, Inches(6.3), Inches(0.95),
             fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
    add_rect(slide, Inches(6.7), tp, Inches(0.07), Inches(0.95), fill_color=color)
    add_text_box(slide, tag,  Inches(6.85), tp + Inches(0.1),
                 Inches(5.8), Inches(0.32), font_size=13, bold=True, color=color)
    add_text_box(slide, desc, Inches(6.85), tp + Inches(0.52),
                 Inches(5.8), Inches(0.28), font_size=11, color=MID_GREY)

# All three tags per team
add_text_box(slide, "Applied to each included team repo (repo-conversion, repo-interfaces, repo-func-config)",
             Inches(6.7), Inches(5.25), Inches(6.3), Inches(0.45),
             font_size=10, italic=True, color=MID_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Key Design Decisions
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Key Design Decisions",
           "Why this strategy was chosen over the alternatives")

decisions = [
    ("Polyrepo, not monorepo",
     "Teams deploy independently. No one team's work blocks another's release. No cross-team merge conflicts."),
    ("manifest.yml as source of truth",
     "CM-owned file in one place. All build decisions are git-committed — full audit trail with who changed what and when."),
    ("included: false — not branch deletion",
     "Deferred team's branch stays intact. No rebasing drama. The branch just isn't wired into the build manifest."),
    ("patch manifest for mid-SIT removal",
     "sprint-12-patch1/manifest.yml preserves the original sprint-12 record. Immutable audit trail for compliance."),
    ("Git tags, not branch names, track environments",
     "Tags survive branch deletion. Ops can always redeploy sit/sprint-12 even after branches are cleaned up."),
    ("Python for YAML report, not bash",
     "Bash grep/awk breaks on inline arrays like [ELM-1201, ELM-1202]. PyYAML handles all valid YAML correctly."),
]

col_w = Inches(4.1)
col_h = Inches(1.6)
gap_x = Inches(0.26)
gap_y = Inches(0.25)
start_x = Inches(0.35)
start_y = Inches(1.35)

for i, (title, body) in enumerate(decisions):
    col = i % 3
    row = i // 3
    lft = start_x + col * (col_w + gap_x)
    tp  = start_y + row * (col_h + gap_y)
    add_rect(slide, lft, tp, col_w, col_h,
             fill_color=WHITE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.75)
    add_rect(slide, lft, tp, col_w, Inches(0.36), fill_color=MID_BLUE)
    add_text_box(slide, title, lft + Inches(0.1), tp + Inches(0.04),
                 col_w - Inches(0.2), Inches(0.28), font_size=11, bold=True, color=WHITE)
    add_text_box(slide, body, lft + Inches(0.12), tp + Inches(0.44),
                 col_w - Inches(0.24), col_h - Inches(0.54),
                 font_size=10.5, color=DARK_GREY, wrap=True)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Live Demo Checklist
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GREY)
header_bar(slide, "Live Demo Checklist",
           "Step-by-step commands to run during the demo")

steps_left = [
    ("1. Show GitHub repos",
     "Open all 5 repos in browser: main branch vs sprint/sprint-12 diff"),
    ("2. Show manifest.yml",
     "cat releases/sprint-12/manifest.yml — highlight included: false for workflow_config"),
    ("3. Generate report",
     "./scripts/generate-report.sh sprint-12\ncat releases/sprint-12/RELEASE-NOTES.md"),
    ("4. Show sprint-12 diff (conversion)",
     "git -C ../repo-conversion diff main sprint/sprint-12"),
    ("5. Show sprint-12 diff (interfaces)",
     "git -C ../repo-interfaces diff main sprint/sprint-12"),
]
steps_right = [
    ("6. Show deferred branch",
     "git -C ../repo-workflow-config log main..sprint/sprint-12 --oneline"),
    ("7. Run tag-environments (SIT)",
     "./scripts/tag-environments.sh sprint-12 sit\ngit -C ../repo-conversion tag -l"),
    ("8. Show no tag on deferred repo",
     "git -C ../repo-workflow-config tag -l\n(empty — correctly excluded)"),
    ("9. Show close-sprint dry run",
     "Open close-sprint.sh and walk through merge + release tag logic"),
    ("10. Show issue template",
     ".github/ISSUE_TEMPLATE/remove-team-mid-sit.md — mid-SIT removal checklist"),
]

for i, (title, cmd) in enumerate(steps_left):
    tp = Inches(1.35 + i * 1.1)
    add_rect(slide, Inches(0.35), tp, Inches(6.3), Inches(0.95),
             fill_color=WHITE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.5)
    add_text_box(slide, title, Inches(0.5), tp + Inches(0.06),
                 Inches(6.0), Inches(0.28), font_size=11, bold=True, color=MID_BLUE)
    add_text_box(slide, cmd, Inches(0.5), tp + Inches(0.42),
                 Inches(6.0), Inches(0.42), font_size=9.5, color=DARK_GREY)

for i, (title, cmd) in enumerate(steps_right):
    tp = Inches(1.35 + i * 1.1)
    add_rect(slide, Inches(6.85), tp, Inches(6.15), Inches(0.95),
             fill_color=WHITE, line_color=RGBColor(0xBB, 0xCF, 0xE0), line_width_pt=0.5)
    add_text_box(slide, title, Inches(7.0), tp + Inches(0.06),
                 Inches(5.9), Inches(0.28), font_size=11, bold=True, color=MID_BLUE)
    add_text_box(slide, cmd, Inches(7.0), tp + Inches(0.42),
                 Inches(5.9), Inches(0.42), font_size=9.5, color=DARK_GREY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary / Next Steps
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BLUE)
add_rect(slide, 0, Inches(2.1), SLIDE_W, Inches(0.06), fill_color=ACCENT_BLUE)

add_text_box(slide, "What This Prototype Demonstrates",
             Inches(1.2), Inches(0.4), Inches(10.9), Inches(0.75),
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

summary = [
    "End-to-end repeatable process from sprint cut to UAT close",
    "CM controls all build decisions via a single manifest.yml — auditable, versioned, CM-gated",
    "Teams work independently with zero cross-repo coordination overhead",
    "Deferred teams never block a release — just set included: false",
    "Environment state is always queryable: which commit is in SIT right now? Check the tag.",
    "ELM traceability end-to-end: work item → feature branch → sprint branch → manifest → release tag",
    "Mid-sprint removals handled cleanly via patch manifests, not branch surgery",
]
text = "\n".join(f"  ▸  {s}" for s in summary)
add_text_box(slide, text,
             Inches(1.0), Inches(2.4), Inches(11.3), Inches(3.5),
             font_size=14, color=RGBColor(0xBD, 0xD7, 0xEE), wrap=True)

add_text_box(slide, "Repositories:",
             Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.32),
             font_size=12, bold=True, color=WHITE)
add_text_box(slide, "github.com/satishreddy13  ·  repo-release-manifest  ·  repo-conversion  ·  repo-interfaces  ·  repo-workflow-config  ·  repo-func-config",
             Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.32),
             font_size=11, color=RGBColor(0x9D, 0xB8, 0xD2), align=PP_ALIGN.CENTER)

footer(slide)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
